from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.services.llm import LLMService
from app.services.text_utils import sentence_split, slugify


class SlideExporter:
    def __init__(self, llm: LLMService, export_dir: Path):
        self.llm = llm
        self.export_dir = export_dir

    def build(self, title: str, chapter: str, context: str) -> Path:
        fallback = self._fallback_outline(title, chapter, context)
        task = """یک طرح اسلاید دانشگاهی تولید کن. خروجی آرایه JSON باشد و هر عضو این فیلدها را داشته باشد:
title, bullets. bullets آرایه‌ای از حداکثر 5 نکته کوتاه باشد. بین 6 تا 12 اسلاید بساز."""
        data, _mode = self.llm.generate_json(task, context, fallback)
        if not isinstance(data, list) or not data:
            data = fallback

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        first = prs.slides.add_slide(prs.slide_layouts[0])
        first.shapes.title.text = title
        first.placeholders[1].text = chapter
        self._format_slide(first)

        for item in data[:15]:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = str(item.get("title", "اسلاید آموزشی"))
            frame = slide.placeholders[1].text_frame
            frame.clear()
            bullets = item.get("bullets", [])
            if not isinstance(bullets, list):
                bullets = [str(bullets)]
            for idx, bullet in enumerate(bullets[:6]):
                paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
                paragraph.text = str(bullet)
                paragraph.level = 0
                paragraph.alignment = PP_ALIGN.RIGHT
                paragraph.font.size = Pt(24)
            self._format_slide(slide)

        filename = f"{slugify(title)}-{slugify(chapter)}.pptx"
        path = self.export_dir / filename
        prs.save(path)
        return path

    @staticmethod
    def _format_slide(slide) -> None:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            for paragraph in shape.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.RIGHT
                for run in paragraph.runs:
                    run.font.name = "Arial"

    @staticmethod
    def _fallback_outline(title: str, chapter: str, context: str) -> list[dict]:
        sentences = sentence_split(context)
        chunks = [sentences[i : i + 4] for i in range(0, min(len(sentences), 32), 4)]
        outline = [{"title": "اهداف یادگیری", "bullets": [f"درک مفاهیم اصلی {chapter}", "مرور نکات کلیدی", "آمادگی برای ارزیابی"]}]
        for index, group in enumerate(chunks[:7], start=1):
            outline.append({"title": f"مفهوم کلیدی {index}", "bullets": group})
        outline.append({"title": "جمع‌بندی", "bullets": [s for s in sentences[:4]]})
        return outline
